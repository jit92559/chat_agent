from graph.state import MainState
from llms.vision_model import get_vision_model

import base64
import pymupdf
import asyncio

from pptx import Presentation
from docx import Document
from dotenv import load_dotenv

from langchain_core.messages import HumanMessage
from fastapi.concurrency import run_in_threadpool

load_dotenv()

vision_model = get_vision_model()


async def describe_image_bytes(
    image_bytes: bytes,
    mime_type: str = "image/png",
) -> str:
    print("i amm doin ocr")
    try:
        image_base64 = base64.b64encode(image_bytes).decode("utf-8")

        message = HumanMessage(
            content=[
                {
                    "type": "text",
                    "text": """
Extract all visible text from this image.
Also describe diagrams, charts, tables, screenshots, UI, and important visual details.
Return clean text suitable for RAG.
""",
                },
                {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:{mime_type};base64,{image_base64}"
                    },
                },
            ]
        )

        response = await vision_model.ainvoke([message])
        return str(response.content)

    except Exception as e:
        return f"[VISION_ERROR] {str(e)}"


async def pdf_extractor_node(state: MainState) -> dict:
    print(" iam at pdf_extrectror node")
    try:
        file_path = state["file_path"]

        def _extract_pdf_sync():
            doc = pymupdf.open(file_path)
            pages_data = []

            for page_no, page in enumerate(doc, start=1):
                page_text = page.get_text("text")

                page_image_bytes = None
                try:
                    pix = page.get_pixmap(matrix=pymupdf.Matrix(2, 2))
                    page_image_bytes = pix.tobytes("png")
                except Exception as e:
                    page_image_bytes = f"[PAGE_OCR_ERROR] {str(e)}"

                embedded_images = []
                images = page.get_images(full=True)

                for image_no, image in enumerate(images, start=1):
                    try:
                        xref = image[0]
                        image_info = doc.extract_image(xref)
                        image_bytes = image_info["image"]

                        if len(image_bytes) < 3000:
                            continue

                        ext = image_info.get("ext", "png").lower()
                        if ext == "jpg":
                            ext = "jpeg"

                        embedded_images.append(
                            {
                                "image_no": image_no,
                                "image_bytes": image_bytes,
                                "mime_type": f"image/{ext}",
                            }
                        )

                    except Exception as e:
                        embedded_images.append(
                            {
                                "image_no": image_no,
                                "error": str(e),
                            }
                        )

                pages_data.append(
                    {
                        "page_no": page_no,
                        "page_text": page_text,
                        "page_image_bytes": page_image_bytes,
                        "embedded_images": embedded_images,
                    }
                )

            doc.close()
            return pages_data

        pages_data = await run_in_threadpool(_extract_pdf_sync)

        all_text = []

        for page in pages_data:
            page_no = page["page_no"]
            page_text = page.get("page_text") or ""

            if page_text.strip():
                all_text.append(
                    f"\n--- Page {page_no} Text ---\n{page_text}"
                )

            page_image_bytes = page.get("page_image_bytes")

            if isinstance(page_image_bytes, bytes):
                page_description = await describe_image_bytes(
                    page_image_bytes,
                    mime_type="image/png",
                )

                if page_description.strip():
                    all_text.append(
                        f"\n--- Page {page_no} Full Page Vision OCR ---\n"
                        f"{page_description}"
                    )

            elif isinstance(page_image_bytes, str):
                all_text.append(
                    f"\n--- Page {page_no} Full Page OCR Error ---\n{page_image_bytes}"
                )

            for image in page.get("embedded_images", []):
                image_no = image.get("image_no")

                if image.get("error"):
                    all_text.append(
                        f"\n--- Page {page_no} Image {image_no} Error ---\n"
                        f"{image.get('error')}"
                    )
                    continue

                image_description = await describe_image_bytes(
                    image["image_bytes"],
                    mime_type=image["mime_type"],
                )

                if image_description.strip():
                    all_text.append(
                        f"\n--- Page {page_no} Image {image_no} Description ---\n"
                        f"{image_description}"
                    )

        extracted_text = "\n".join(all_text)

        if not extracted_text.strip():
            return {
                "extracted_text": "",
                "status": "failed",
                "error": "No text extracted from PDF",
            }

        return {
            "extracted_text": extracted_text,
            "status": "pdf_extracted",
            "error": None,
        }

    except Exception as e:
        return {
            "extracted_text": "",
            "status": "failed",
            "error": str(e),
        }


async def ppt_extractor_node(state: MainState) -> dict:
    try:
        file_path = state["file_path"]

        def _extract_ppt_sync():
            prs = Presentation(file_path)
            slides_data = []

            for slide_no, slide in enumerate(prs.slides, start=1):
                slide_text = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                images_data = []

                for shape_no, shape in enumerate(slide.shapes, start=1):
                    try:
                        if hasattr(shape, "image"):
                            image_bytes = shape.image.blob

                            if len(image_bytes) < 3000:
                                continue

                            ext = shape.image.ext or "png"
                            ext = ext.lower()

                            if ext == "jpg":
                                ext = "jpeg"

                            images_data.append(
                                {
                                    "shape_no": shape_no,
                                    "image_bytes": image_bytes,
                                    "mime_type": f"image/{ext}",
                                }
                            )

                    except Exception as e:
                        images_data.append(
                            {
                                "shape_no": shape_no,
                                "error": str(e),
                            }
                        )

                slides_data.append(
                    {
                        "slide_no": slide_no,
                        "slide_text": slide_text,
                        "images_data": images_data,
                    }
                )

            return slides_data

        slides_data = await run_in_threadpool(_extract_ppt_sync)

        all_text = []

        for slide in slides_data:
            slide_no = slide["slide_no"]
            slide_text = slide.get("slide_text") or []

            if slide_text:
                all_text.append(
                    f"\n--- Slide {slide_no} Text ---\n"
                    + "\n".join(slide_text)
                )

            for image in slide.get("images_data", []):
                shape_no = image.get("shape_no")

                if image.get("error"):
                    all_text.append(
                        f"\n--- Slide {slide_no} Image {shape_no} Error ---\n"
                        f"{image.get('error')}"
                    )
                    continue

                image_description = await describe_image_bytes(
                    image["image_bytes"],
                    mime_type=image["mime_type"],
                )

                if image_description.strip():
                    all_text.append(
                        f"\n--- Slide {slide_no} Image {shape_no} Description ---\n"
                        f"{image_description}"
                    )

        extracted_text = "\n".join(all_text)

        if not extracted_text.strip():
            return {
                "extracted_text": "",
                "status": "failed",
                "error": "No text extracted from PPTX",
            }

        return {
            "extracted_text": extracted_text,
            "status": "ppt_extracted",
            "error": None,
        }

    except Exception as e:
        return {
            "extracted_text": "",
            "status": "failed",
            "error": str(e),
        }


async def image_extractor_node(state: MainState) -> dict:
    try:
        file_path = state["file_path"]

        def _read_image_sync():
            with open(file_path, "rb") as f:
                return f.read()

        image_bytes = await run_in_threadpool(_read_image_sync)

        description = await describe_image_bytes(image_bytes)

        if not description.strip():
            return {
                "extracted_text": "",
                "status": "failed",
                "error": "No text extracted from image",
            }

        return {
            "extracted_text": description,
            "status": "image_extracted",
            "error": None,
        }

    except Exception as e:
        return {
            "extracted_text": "",
            "status": "failed",
            "error": str(e),
        }


async def doc_extractor_node(state: MainState) -> dict:
    try:
        file_path = state["file_path"]

        def _extract_doc_sync():
            doc = Document(file_path)
            return "\n".join(
                para.text for para in doc.paragraphs if para.text.strip()
            )

        text = await run_in_threadpool(_extract_doc_sync)

        if not text.strip():
            return {
                "extracted_text": "",
                "status": "failed",
                "error": "No text extracted from DOCX",
            }

        return {
            "extracted_text": text,
            "status": "doc_extracted",
            "error": None,
        }

    except Exception as e:
        return {
            "extracted_text": "",
            "status": "failed",
            "error": str(e),
        }


async def txt_extractor_node(state: MainState) -> dict:
    try:
        file_path = state["file_path"]

        def _read_txt_sync():
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

        text = await run_in_threadpool(_read_txt_sync)

        if not text.strip():
            return {
                "extracted_text": "",
                "status": "failed",
                "error": "No text extracted from TXT",
            }

        return {
            "extracted_text": text,
            "status": "txt_extracted",
            "error": None,
        }

    except Exception as e:
        return {
            "extracted_text": "",
            "status": "failed",
            "error": str(e),
        }


def error_node(state: MainState) -> dict:
    return {
        "error": f"Error: {state.get('error', 'Unknown error occurred')}",
        "status": "failed",
    }